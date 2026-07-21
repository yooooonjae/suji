"""포트폴리오 PPT 생성 — 「수지(收支)」 프로젝트 소개 (이미지 포함, 16:9).

실행: venv/bin/python src/build/portfolio_ppt.py
산출: 포트폴리오/수지_부동산수지분석_포트폴리오.pptx
사전 준비: 포트폴리오/이미지/01~05 스크린샷 (헤드리스 크롬, 주황 테마)
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from PIL import Image

ROOT = Path(__file__).resolve().parents[2]
IMG = ROOT / "포트폴리오" / "이미지"
OUT = ROOT / "포트폴리오" / "수지_부동산수지분석_포트폴리오.pptx"

# 브랜드 토큰 (사이트 주황 제도지 테마와 동일 계열)
INK = RGBColor(0x2B, 0x20, 0x16)     # 웜 먹색
INK2 = RGBColor(0x6B, 0x5A, 0x49)    # 웜 회갈
ORANGE = RGBColor(0xC2, 0x41, 0x0C)  # 제도펜 주황 (--blueprint)
SEAL = RGBColor(0x1E, 0x5D, 0x95)    # 인장 청람 (반전 대비색)
PAPER = RGBColor(0xFA, 0xF7, 0xF2)   # 웜 제도지
CARD = RGBColor(0xFF, 0xFD, 0xFA)

SW, SH = Inches(13.333), Inches(7.5)  # 16:9


def base(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = PAPER
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT):
    """runs: [(문자열, 크기pt, 볼드, 색)] — 줄 단위."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (t, size, bold, color) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Apple SD Gothic Neo"
    return tb


def rule(slide, x, y, w, color=ORANGE, h=Pt(2.2)):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def pic_fit(slide, path, x, y, max_w, max_h, top_crop_ratio=None):
    """이미지를 비율 유지로 박스 안에 배치. top_crop_ratio 지정 시 상단부만 사용."""
    p = Path(path)
    if top_crop_ratio:
        im = Image.open(p)
        im = im.crop((0, 0, im.width, int(im.height * top_crop_ratio)))
        tmp = p.parent / f"_crop_{p.name}"
        im.save(tmp)
        p = tmp
    im = Image.open(p)
    ratio = im.width / im.height
    w, h = max_w, Emu(int(max_w / ratio))
    if h > max_h:
        h = max_h
        w = Emu(int(max_h * ratio))
    return slide.shapes.add_picture(str(p), x + (max_w - w) // 2, y, w, h)


def header(slide, no, title, sub):
    text(slide, Inches(0.7), Inches(0.42), Inches(1.0), Inches(0.8),
         [(no, 30, True, ORANGE)])
    text(slide, Inches(1.5), Inches(0.42), Inches(10.5), Inches(0.6),
         [(title, 26, True, INK)])
    text(slide, Inches(1.5), Inches(1.02), Inches(11), Inches(0.4),
         [(sub, 13, False, INK2)])
    rule(slide, Inches(0.7), Inches(1.5), Inches(11.9))


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH

    # ── 1. 표지 ──────────────────────────────────────
    s = base(prs)
    pic_fit(s, IMG / "01-표지.png", Inches(0), Inches(0), SW, SH)
    from pptx.enum.shapes import MSO_SHAPE
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.7), SW, Inches(0.8))
    bar.fill.solid(); bar.fill.fore_color.rgb = INK; bar.line.fill.background()
    text(s, Inches(0.7), Inches(6.86), Inches(12), Inches(0.5),
         [("수지(收支) — 대한민국 부동산 개발의 손익 구조  ·  개인 연구 포트폴리오  ·  서윤재  ·  yoonjae.pages.dev", 14.5, True, PAPER)])

    # ── 2. 개요 ──────────────────────────────────────
    s = base(prs)
    header(s, "01", "프로젝트 개요", "공공 데이터 → 재무 모델 → 인터랙티브 웹 리포트 — 매월 자동 갱신되는 연구 플랫폼")
    text(s, Inches(0.7), Inches(1.9), Inches(5.9), Inches(4.6), [
        ("무엇을 만들었나", 17, True, ORANGE),
        ("대한민국 부동산 개발사업(신축분양·재개발·재건축·리모델링)의", 14, False, INK),
        ("수지분석을 데이터와 검증된 모델로 수행하는 연구 웹 리포트.", 14, False, INK),
        ("", 8, False, INK),
        ("누구를 위해", 17, True, ORANGE),
        ("시행사 · 자산운용사 · 컨설팅 — 딜 스크리닝의 초기 검토 흐름을", 14, False, INK),
        ("재현하는 정밀도와 기관 보고서 격의 완성도를 지향.", 14, False, INK),
        ("", 8, False, INK),
        ("무엇이 다른가", 17, True, ORANGE),
        ("① 토지 용도지역→용적률→세대수까지 법령 기반 자동 도출", 14, False, INK),
        ("② 정비 비례율·분담금, 수익형 NOI·환원율까지 실시간 계산기", 14, False, INK),
        ("③ 딥러닝·파운데이션 모델까지 겨룬 예측 벤치마크의 정직한 공개", 14, False, INK),
    ])
    text(s, Inches(7.0), Inches(1.9), Inches(5.6), Inches(4.8), [
        ("한눈에 보는 규모", 17, True, ORANGE),
        ("", 5, False, INK),
        ("236,388건   아파트 실거래 (매매+분양권)", 15.5, True, INK),
        ("195,863건   서울 25개 구 아파트 전수", 15.5, True, INK),
        ("153,285건   상업·업무·오피스텔·토지 실거래", 15.5, True, INK),
        ("77계열      시장 시계열 (월별 10년)", 15.5, True, INK),
        ("14종        공공 데이터셋 (7개 기관 API)", 15.5, True, INK),
        ("6모델       예측 벤치마크", 15.5, True, INK),
        ("162개       자동화 검증 테스트", 15.5, True, INK),
        ("", 5, False, INK),
        ("한국부동산원 R-ONE · 국토부 RTMS·건축HUB · KOSIS ·", 13, False, INK2),
        ("한국은행 ECOS · HUG · 소상공인시장진흥공단", 13, False, INK2),
    ])

    # ── 3. 파이프라인 ─────────────────────────────────
    s = base(prs)
    header(s, "02", "데이터 파이프라인", "수집(14종) → 검증 → 모델 → 예측 → 웹·배포 — launchd 매월 자동 갱신, 실패는 격리·통보")
    steps = [
        ("수집", "10개 수집기\n원본 전량 캐시\n백오프·폴백 명시"),
        ("검증", "단위 통일(원·㎡)\n물리 범위 게이트\n결측은 보간 없이 기록"),
        ("모델", "수지 4모드+수익형\n용도지역 모듈\nPython↔JS 패리티"),
        ("예측", "SARIMA·LGBM·LSTM\nChronos 제로샷\n분기마다 재학습"),
        ("웹·배포", "8챕터+EDA 리포트\n실시간 계산기\n월간 자동 배포"),
    ]
    x = Inches(0.7)
    for i, (t, d) in enumerate(steps):
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.2), Inches(2.2), Inches(2.6))
        box.fill.solid(); box.fill.fore_color.rgb = CARD
        box.line.color.rgb = ORANGE; box.line.width = Pt(1.2)
        tf = box.text_frame; tf.word_wrap = True
        p0 = tf.paragraphs[0]; r0 = p0.add_run(); r0.text = t
        r0.font.size = Pt(17); r0.font.bold = True; r0.font.color.rgb = ORANGE
        for line in d.split("\n"):
            p = tf.add_paragraph(); r = p.add_run(); r.text = line
            r.font.size = Pt(12); r.font.color.rgb = INK2
        if i < 4:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.24), Inches(3.25), Inches(0.34), Inches(0.35))
            ar.fill.solid(); ar.fill.fore_color.rgb = ORANGE; ar.line.fill.background()
        x += Inches(2.6)
    text(s, Inches(0.7), Inches(5.3), Inches(12), Inches(1.6), [
        ("운영 원칙 — 실무 시스템 운영에서 얻은 교훈의 이식", 15, True, INK),
        ("API 오류·쿼터를 침묵으로 삼키지 않는다 · 모든 집계는 원본 대조 검증 스크립트를 동반한다 ·", 13, False, INK2),
        ("빌드는 원자적 스왑(실패 시 직전 사이트 보존) · 행정구역 개편 같은 원천 변화도 수집기가 흡수한다", 13, False, INK2),
    ])

    # ── 4. 수지모델 ───────────────────────────────────
    s = base(prs)
    header(s, "03", "수지분석 모델 — 4모드 + 수익형, 이중 구현, 교차 검증", "신축분양 · 재개발 · 재건축 · 리모델링 · 오피스/상업 수익형")
    text(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(4.8), [
        ("구조", 16, True, ORANGE),
        ("수입(분양) − 지출 여섯 항목(토지·공사·간접·판매·금융·예비)", 14, False, INK),
        ("→ 분기 현금흐름 → 개발이익·마진·ROE·NPV·IRR", 14, False, INK),
        ("", 8, False, INK),
        ("정비사업 고유 계산", 16, True, ORANGE),
        ("비례율 = (총수입−총사업비) ÷ 종전자산평가액", 14, False, INK),
        ("분담금 = 조합원분양가 − 권리가액 · 임대의무·현금청산 반영", 14, False, INK),
        ("", 8, False, INK),
        ("수익형(오피스·상업)", 16, True, ORANGE),
        ("NOI = 임대면적×임대료×(1−공실)×(1−경비율)", 14, False, INK),
        ("매각가치 = NOI ÷ 환원율 — 준공 시점 일시 유입으로 편성", 14, False, INK),
    ])
    text(s, Inches(7.1), Inches(1.9), Inches(5.5), Inches(4.8), [
        ("검증 체계", 16, True, SEAL),
        ("", 6, False, INK),
        ("162개 pytest — 수기검산 대조·경계·회귀", 14, True, INK),
        ("", 4, False, INK),
        ("Python ↔ JavaScript 패리티 78조", 14, True, INK),
        ("전 출력 |Δ| < 1e-6 강제 (실측 최대 상대오차 2.5e-15)", 13, False, INK2),
        ("", 4, False, INK),
        ("독립 재구현 3-way 대조", 14, True, INK),
        ("명세만 보고 새로 짠 검산기와 양 엔진의 수치 완전 일치", 13, False, INK2),
        ("", 4, False, INK),
        ("적대적 교차 검증 4회", 14, True, INK),
        ("외부 모델·다중 에이전트 — 수식·이식·접근성·UX 실측", 13, False, INK2),
    ])

    # ── 5. 예측 벤치마크 ──────────────────────────────
    s = base(prs)
    header(s, "04", "예측 벤치마크 — 정직한 결과", "시도별 매매가격지수 · 롤링 오리진 백테스트 · 동일 조건 6모델")
    rows = [("1", "SARIMA (statsmodels)", "0.680", "우승 — 소표본의 고전 강자"),
            ("2", "Chronos-Bolt (HuggingFace 제로샷)", "0.922", "학습 없이 2위 — 파운데이션 모델의 저력"),
            ("3", "Naive (마지막값)", "1.088", "이기기 어려운 기준선"),
            ("4", "LightGBM (분위수·라그 피처)", "1.614", ""),
            ("5", "계절 Naive", "1.869", ""),
            ("6", "LSTM (PyTorch, 시도 풀링)", "3.046", "소표본에서 딥러닝의 한계를 그대로 보고")]
    y = Inches(2.0)
    for rank, name, mae, note in rows:
        c = ORANGE if rank == "1" else (INK if rank in "23" else INK2)
        text(s, Inches(0.8), y, Inches(0.6), Inches(0.4), [(rank, 16, rank == "1", c)])
        text(s, Inches(1.5), y, Inches(5.6), Inches(0.4), [(name, 15, rank == "1", c)])
        text(s, Inches(7.2), y, Inches(1.4), Inches(0.4), [("MAE " + mae, 15, rank == "1", c)])
        text(s, Inches(8.8), y, Inches(4.2), Inches(0.4), [(note, 12.5, False, INK2)])
        y += Inches(0.52)
    text(s, Inches(0.8), y + Inches(0.25), Inches(11.8), Inches(1.4), [
        ("정직한 보고 원칙", 15, True, SEAL),
        ("120개월 소표본에서 딥러닝이 지는 것은 예상된 결과이며, 이를 감추지 않고 벤치마크 표로 공개한다.", 13.5, False, INK),
        ("미래 정보 누출 차단(라그 피처만) · 시드 고정 재현성 · 80% 분위수 구간으로 불확실성 명시.", 13.5, False, INK),
    ])

    # ── 6~8. 웹 쇼케이스 ─────────────────────────────
    for no, title, sub, img, crop in [
        ("05", "웹 — 시장 챕터", "17개 시도 스몰멀티플(선택 연동) · 서울 25개 구 전수 195,863건 — 강남3구가 비강남의 1.87배", IMG / "02-시장.png", 0.55),
        ("06", "웹 — 실시간 수지 계산기", "토지(용도지역)→세대수 자동 도출 · 4모드+수익형 · 워터폴·게이지 · 시나리오 A/B", IMG / "03-계산기.png", 0.55),
        ("07", "웹 — 민감도 & 다크 모드", "토네이도 · 손익분기 히트맵(셀 수치 직접 표기, 계산기 연동) · 완전한 다크 테마 · 모바일 실측 검증", IMG / "04-민감도상권.png", 0.50),
    ]:
        s = base(prs)
        header(s, no, title, sub)
        pic_fit(s, img, Inches(0.7), Inches(1.75), Inches(11.9), Inches(5.4), top_crop_ratio=crop)
    s = prs.slides[-1]
    pic_fit(s, IMG / "05-다크.png", Inches(7.4), Inches(3.9), Inches(5.3), Inches(3.2))

    # ── 9. 기술 스택·마무리 ───────────────────────────
    s = base(prs)
    header(s, "08", "기술 스택 & 이 프로젝트가 증명하는 것", "")
    text(s, Inches(0.7), Inches(1.9), Inches(5.9), Inches(4.8), [
        ("스택", 16, True, ORANGE),
        ("Python (수집·모델·빌드) · pytest 162", 14, False, INK),
        ("PyTorch · LightGBM · statsmodels · HuggingFace Chronos", 14, False, INK),
        ("바닐라 JS/SVG 차트 자작 — 드래그 확대·월/분기/연 전환, 의존 0", 14, False, INK),
        ("공공 API 7개 기관 14종 (data.go.kr · ECOS · KOSIS · R-ONE 등)", 14, False, INK),
        ("", 8, False, INK),
        ("디자인", 16, True, ORANGE),
        ("“웜 제도지 + 주황 제도펜” 에디토리얼 — 도곽·치수선·인장 모티프", 14, False, INK),
        ("검증기 통과 팔레트 · 라이트/다크 · 모바일 390px 실측 검증", 14, False, INK),
    ])
    text(s, Inches(7.0), Inches(1.9), Inches(5.6), Inches(4.8), [
        ("증명하는 역량", 16, True, SEAL),
        ("", 6, False, INK),
        ("도메인 — 개발·정비사업 수지 구조, 법령(국계법·주택법·도정법)", 14, False, INK),
        ("데이터 — 공공 API 전 주기(발급·수집·검증·재현) 운영", 14, False, INK),
        ("모델링 — 재무 모델 TDD, 시계열 ML/DL 벤치마크 설계", 14, False, INK),
        ("엔지니어링 — 이중 구현 패리티, 파이프라인 멱등성·자동화", 14, False, INK),
        ("프로덕트 — 기관 독자를 위한 서사·시각화·접근성 설계", 14, False, INK),
        ("", 10, False, INK),
        ("공개: https://yoonjae.pages.dev — 매월 자동 갱신", 13.5, True, ORANGE),
    ])

    OUT.parent.mkdir(exist_ok=True)
    prs.save(OUT)
    for f in IMG.glob("_crop_*"):
        f.unlink()
    print(f"PPT 저장: {OUT} ({OUT.stat().st_size/1024:.0f} KB)")


if __name__ == "__main__":
    build()
